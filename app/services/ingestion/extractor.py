import io


def extract_text(content: bytes, filename: str) -> str:
    """Extract plain text from PDF, Word, Excel, PowerPoint, CSV or text bytes."""
    name = filename.lower()
    if name.endswith(".pdf"):
        return _extract_pdf(content)
    if name.endswith(".docx"):
        return _extract_docx(content)
    if name.endswith(".xlsx"):
        return _extract_xlsx(content)
    if name.endswith(".pptx"):
        return _extract_pptx(content)
    if name.endswith(".csv"):
        return _extract_csv(content)
    return content.decode("utf-8", errors="replace")


def _extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader  # type: ignore[import-untyped]

    reader = PdfReader(io.BytesIO(content))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages)


def _extract_docx(content: bytes) -> str:
    from docx import Document  # type: ignore[import-untyped]

    doc = Document(io.BytesIO(content))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    # Include table cell text (contracts/invoices often use tables).
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _extract_xlsx(content: bytes) -> str:
    from openpyxl import load_workbook  # type: ignore[import-untyped]

    # read_only=False required to access number_format on each cell so we can
    # reconstruct RUTs stored as formatted integers (e.g. 123456789 → 12.345.678-9).
    wb = load_workbook(io.BytesIO(content), data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows():
            cells = [s for c in row if (s := _xlsx_cell(c))]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"=== Hoja: {sheet.title} ===\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts)


def _xlsx_cell(cell: object) -> str:
    """Return a cell's display string, reconstructing RUT-formatted numbers.

    When openpyxl reads an integer whose number_format contains dots and a dash
    (e.g. "##.###.###-0"), the raw value is e.g. 123456789 and the scrubber
    regex would miss it. We detect that pattern and rebuild "12.345.678-9".
    """
    val = getattr(cell, "value", None)
    if val is None:
        return ""
    if isinstance(val, int):
        fmt: str = getattr(cell, "number_format", "") or ""
        if "." in fmt and "-" in fmt:
            s = str(abs(val))
            if 7 <= len(s) <= 9:
                body, verif = s[:-1], s[-1]
                parts: list[str] = []
                while body:
                    parts.append(body[-3:] if len(body) >= 3 else body)
                    body = body[:-3]
                return ".".join(reversed(parts)) + f"-{verif}"
    return str(val).strip()


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation  # type: ignore[import-untyped]

    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        lines.append(" | ".join(cells))
        if lines:
            parts.append(f"=== Diapositiva {i} ===\n" + "\n".join(lines))
    return "\n\n".join(parts)


def _extract_csv(content: bytes) -> str:
    import csv

    # utf-8-sig strips a BOM if present; Chilean CSVs often use ';' as delimiter
    # because ',' is the decimal separator, so we sniff it.
    text = content.decode("utf-8-sig", errors="replace")
    try:
        dialect = csv.Sniffer().sniff(text[:4096], delimiters=",;\t|")
    except csv.Error:
        dialect = csv.excel
    rows: list[str] = []
    for row in csv.reader(io.StringIO(text), dialect):
        cells = [c.strip() for c in row if c.strip()]
        if cells:
            rows.append(" | ".join(cells))
    return "\n".join(rows)
